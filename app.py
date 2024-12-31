import streamlit as st
import tensorflow as tf
from tensorflow.keras.preprocessing import image
import numpy as np
from PIL import Image
import fitz  # For PDF handling (PyMuPDF)
from pptx import Presentation
import docx
import io

# === Load the Pre-trained Model ===
MODEL_PATH = "sucker_rod_pump_model.h5"
IMG_HEIGHT = 224
IMG_WIDTH = 224
model = tf.keras.models.load_model(MODEL_PATH)

# === Class Labels and Solutions ===
class_indices = {
    0: "Ideal Card",
    1: "Gas Interference",
    2: "Fluid Pound",
    3: "Rod Parted",
    4: "Pump Hitting Up",
    5: "Pump Hitting Down",
    6: "Bent Barrel",
    7: "Tubing Movement",
    8: "Worn or Split Barrel",
    9: "Worn Plunger or Traveling Valve",
    10: "Worn Standing Valve"
}

solutions = {
    "Gas Interference": """**Causes**: 
    - Incomplete pump filling, reduced fluid load, and abnormal pump behavior.
    **Solutions**: 
    - Optimize Pump Settings: Reduce stroke per minute (SPM) and increase stroke length.
    - Enhance Gas Separation: Install downhole gas separators and improve casing gas venting.
    - Maintain Backpressure: Use surface backpressure valves.
    - Use Chemical Treatments: Inject foam breakers or surfactants to minimize gas impact.
    - Adjust Well Configuration: Use tubing anchors and optimize perforation placement.
    - Long-Term Fixes: Install ESPs or gas lift systems.""",
    "Fluid Pound": """**Causes**: 
    - Liquid fallback due to incomplete filling or over-pumping.
    **Solutions**:
    - Check pump inlet for obstructions.
    - Reduce SPM and adjust pump stroke settings.
    - Verify fluid levels in the well.""",
    "Rod Parted": """**Causes**: 
    - Excessive tensile stress or corrosion.
    **Solutions**:
    - Conduct rod inspection and use higher-quality rods.
    - Check alignment and fluid loading conditions.""",
    "Pump Hitting Up": """**Causes**: 
    - Insufficient plunger clearance at the top.
    **Solutions**:
    - Adjust stroke settings to prevent plunger impact at the top.""",
    "Pump Hitting Down": """**Causes**: 
    - Insufficient clearance at the bottom.
    **Solutions**:
    - Adjust stroke settings to prevent plunger impact at the bottom.""",
    "Bent Barrel": """**Causes**: 
    - Mechanical damage or deformation of the barrel.
    **Solutions**:
    - Replace the bent barrel.
    - Use centralizers to prevent misalignment.""",
    "Tubing Movement": """**Causes**: 
    - Tubing instability due to improper anchoring.
    **Solutions**:
    - Secure tubing with anchors.
    - Inspect well configuration.""",
    "Worn or Split Barrel": """**Causes**: 
    - Abrasive wear or splitting from mechanical stress.
    **Solutions**:
    - Replace the barrel.
    - Use protective coatings to prevent wear.""",
    "Worn Plunger or Traveling Valve": """**Causes**: 
    - Excessive wear from abrasion or corrosion.
    **Solutions**:
    - Replace the plunger or traveling valve.
    - Use wear-resistant materials.""",
    "Worn Standing Valve": """**Causes**: 
    - Erosion or wear from fluid dynamics.
    **Solutions**:
    - Replace the standing valve.
    - Use better quality materials.""",
}

# === Image Preprocessing Function ===
def preprocess_image(img):
    img = img.resize((IMG_HEIGHT, IMG_WIDTH))
    img_array = np.expand_dims(np.array(img) / 255.0, axis=0)
    return img_array

# === Extract Images from Documents ===
def extract_images_from_pdf(pdf_file):
    doc = fitz.open(pdf_file)
    images = []
    for page in doc:
        for img in page.get_images(full=True):
            xref = img[0]
            base_image = doc.extract_image(xref)
            images.append(Image.open(io.BytesIO(base_image["image"])))
    return images

def extract_images_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Type 13 is Picture
                image_stream = shape.image.blob
                images.append(Image.open(io.BytesIO(image_stream)))
    return images

def extract_images_from_docx(docx_file):
    doc = docx.Document(docx_file)
    images = []
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            img = doc.part.related_parts[rel.target_ref]
            images.append(Image.open(io.BytesIO(img.blob)))
    return images

# === Classification Function ===
def classify_image(img):
    img_array = preprocess_image(img)
    predictions = model.predict(img_array)
    predicted_class = class_indices[np.argmax(predictions)]
    confidence = np.max(predictions) * 100
    return predicted_class, confidence

# === Streamlit App ===
st.title("Sucker Rod Pump Dynacard Classification")

uploaded_file = st.file_uploader("Upload an image, PDF, PowerPoint, or Word document", type=["jpg", "png", "pdf", "pptx", "docx"])

if uploaded_file:
    if uploaded_file.name.endswith(".pdf"):
        extracted_images = extract_images_from_pdf(uploaded_file)
    elif uploaded_file.name.endswith(".pptx"):
        extracted_images = extract_images_from_pptx(uploaded_file)
    elif uploaded_file.name.endswith(".docx"):
        extracted_images = extract_images_from_docx(uploaded_file)
    else:
        extracted_images = [Image.open(uploaded_file)]

    for img in extracted_images:
        st.image(img, caption="Uploaded Image", use_column_width=True)
        predicted_class, confidence = classify_image(img)
        st.write(f"**Prediction**: {predicted_class} (Confidence: {confidence:.2f}%)")
        if predicted_class != "Ideal Card":
            st.write(solutions.get(predicted_class, "No specific solution available."))
